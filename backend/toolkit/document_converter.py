"""
Document to PDF Conversion Service
Converts Word, Excel, PowerPoint documents to PDF format using professional libraries.
"""

import os
import io
import logging
import tempfile
import subprocess
import shutil
from pathlib import Path
from typing import Tuple, Optional

# Professional document conversion
try:
    from docx2pdf import convert as docx2pdf_convert
    HAS_DOCX2PDF = True
except ImportError:
    HAS_DOCX2PDF = False

# Excel conversion
try:
    from openpyxl import load_workbook
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False

# PowerPoint conversion
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# PDF generation
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.lib import colors
from PIL import Image as PILImage

logger = logging.getLogger(__name__)

SUPPORTED_FORMATS = {
    'docx': 'Word Document',
    'doc': 'Word Document',
    'xlsx': 'Excel Spreadsheet',
    'xls': 'Excel Spreadsheet',
    'pptx': 'PowerPoint Presentation',
    'ppt': 'PowerPoint Presentation',
    'txt': 'Text Document',
    'png': 'Image',
    'jpg': 'Image',
    'jpeg': 'Image',
}

MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB


class DocumentConverter:
    """Convert various document formats to PDF"""
    
    @staticmethod
    def validate_file(file_path: str, file_size: int) -> Tuple[bool, str]:
        """Validate file format and size"""
        if file_size > MAX_FILE_SIZE:
            return False, f"File size exceeds {MAX_FILE_SIZE / 1024 / 1024}MB limit"
        
        ext = Path(file_path).suffix.lower().lstrip('.')
        if ext not in SUPPORTED_FORMATS:
            return False, f"Unsupported format: {ext}"
        
        return True, "Valid"
    
    @staticmethod
    def docx_to_pdf(file_content: bytes) -> Optional[bytes]:
        """Convert DOCX to PDF using Microsoft Word via docx2pdf"""
        if not HAS_DOCX2PDF:
            logger.error("docx2pdf library is not installed")
            return None
            
        try:
            # Create temporary files
            with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp_docx:
                tmp_docx.write(file_content)
                tmp_docx_path = tmp_docx.name
            
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_pdf:
                tmp_pdf_path = tmp_pdf.name
            
            try:
                # Convert using Microsoft Word - preserves ALL formatting
                docx2pdf_convert(tmp_docx_path, tmp_pdf_path)
                
                # Read the generated PDF
                with open(tmp_pdf_path, 'rb') as pdf_file:
                    pdf_bytes = pdf_file.read()
                
                return pdf_bytes
                
            finally:
                # Clean up
                try:
                    os.unlink(tmp_docx_path)
                    os.unlink(tmp_pdf_path)
                except:
                    pass
                    
        except Exception as e:
            logger.error(f"Error converting DOCX: {e}", exc_info=True)
            return None
    
    @staticmethod
    def xlsx_to_pdf(file_content: bytes) -> Optional[bytes]:
        """Convert XLSX to PDF with 3-tier fallback: LibreOffice → MS Office → Basic extraction"""
        # Tier 1: Try LibreOffice (best quality, cross-platform)
        result = DocumentConverter._convert_with_libreoffice(file_content, '.xlsx')
        if result is not None:
            logger.info("XLSX converted successfully with LibreOffice")
            return result
        
        # Tier 2: Try Microsoft Office (Windows only, requires Excel installed)
        logger.warning("LibreOffice not available, trying Microsoft Office")
        result = DocumentConverter._xlsx_with_msoffice(file_content)
        if result is not None:
            logger.info("XLSX converted successfully with Microsoft Office")
            return result
        
        # Tier 3: Fall back to basic table extraction
        logger.warning("Microsoft Office not available, using basic extraction")
        return DocumentConverter._xlsx_table_fallback(file_content)
    
    @staticmethod
    def _xlsx_with_msoffice(file_content: bytes) -> Optional[bytes]:
        """Convert XLSX using Microsoft Excel COM automation (Windows only)"""
        try:
            import comtypes.client
        except ImportError:
            logger.debug("comtypes not available for Microsoft Office automation")
            return None
        
        try:
            # Create temporary files
            with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp_xlsx:
                tmp_xlsx.write(file_content)
                tmp_xlsx_path = os.path.abspath(tmp_xlsx.name)
            
            tmp_pdf_path = os.path.abspath(tmp_xlsx_path.replace('.xlsx', '.pdf'))
            
            try:
                # Initialize Excel COM object
                excel = comtypes.client.CreateObject("Excel.Application")
                excel.Visible = False
                excel.DisplayAlerts = False
                
                # Open workbook
                workbook = excel.Workbooks.Open(tmp_xlsx_path)
                
                # Save as PDF (format 57 = PDF)
                workbook.ExportAsFixedFormat(0, tmp_pdf_path)
                
                # Close workbook
                workbook.Close(False)
                excel.Quit()
                
                # Read the PDF
                with open(tmp_pdf_path, 'rb') as pdf_file:
                    return pdf_file.read()
                    
            finally:
                # Clean up
                try:
                    os.unlink(tmp_xlsx_path)
                    os.unlink(tmp_pdf_path)
                except:
                    pass
                    
        except Exception as e:
            logger.debug(f"Microsoft Excel conversion failed: {e}")
            return None
    
    @staticmethod
    def _xlsx_table_fallback(file_content: bytes) -> Optional[bytes]:
        """Fallback: Extract tables from XLSX when LibreOffice is unavailable"""
        if not HAS_OPENPYXL:
            return None
        
        try:
            wb = load_workbook(io.BytesIO(file_content))
            pdf_buffer = io.BytesIO()
            doc = SimpleDocTemplate(pdf_buffer, pagesize=letter)
            styles = getSampleStyleSheet()
            story = []
            
            for sheet_idx, sheet_name in enumerate(wb.sheetnames):
                ws = wb[sheet_name]
                if sheet_idx > 0:
                    story.append(PageBreak())
                story.append(Paragraph(f"<b>{sheet_name}</b>", styles['Heading2']))
                story.append(Spacer(1, 0.2 * inch))
                
                table_data = []
                for row in ws.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        table_data.append([str(cell) if cell is not None else '' for cell in row])
                
                if table_data:
                    t = Table(table_data)
                    t.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                        ('FONTSIZE', (0, 0), (-1, 0), 10),
                        ('GRID', (0, 0), (-1, -1), 1, colors.black),
                    ]))
                    story.append(t)
            
            if not story:
                story.append(Paragraph("Empty spreadsheet", styles['Normal']))
            
            doc.build(story)
            pdf_buffer.seek(0)
            return pdf_buffer.getvalue()
        except Exception as e:
            logger.error(f"Fallback XLSX conversion failed: {e}", exc_info=True)
            return None
    
    @staticmethod
    def pptx_to_pdf(file_content: bytes) -> Optional[bytes]:
        """Convert PPTX to PDF with 3-tier fallback: LibreOffice → MS Office → Basic extraction"""
        # Tier 1: Try LibreOffice (best quality, cross-platform)
        result = DocumentConverter._convert_with_libreoffice(file_content, '.pptx')
        if result is not None:
            logger.info("PPTX converted successfully with LibreOffice")
            return result
        
        # Tier 2: Try Microsoft Office (Windows only, requires PowerPoint installed)
        logger.warning("LibreOffice not available, trying Microsoft Office")
        result = DocumentConverter._pptx_with_msoffice(file_content)
        if result is not None:
            logger.info("PPTX converted successfully with Microsoft Office")
            return result
        
        # Tier 3: Fall back to basic text extraction
        logger.warning("Microsoft Office not available, using basic extraction")
        return DocumentConverter._pptx_text_fallback(file_content)
    
    @staticmethod
    def _pptx_with_msoffice(file_content: bytes) -> Optional[bytes]:
        """Convert PPTX using Microsoft PowerPoint COM automation (Windows only)"""
        try:
            import comtypes.client
        except ImportError:
            logger.debug("comtypes not available for Microsoft Office automation")
            return None
        
        try:
            # Create temporary files
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_pptx:
                tmp_pptx.write(file_content)
                tmp_pptx_path = os.path.abspath(tmp_pptx.name)
            
            tmp_pdf_path = os.path.abspath(tmp_pptx_path.replace('.pptx', '.pdf'))
            
            try:
                # Initialize PowerPoint COM object
                powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
                powerpoint.Visible = 1
                
                # Open presentation
                presentation = powerpoint.Presentations.Open(tmp_pptx_path, WithWindow=False)
                
                # Save as PDF (format 32 = PDF)
                presentation.SaveAs(tmp_pdf_path, 32)
                
                # Close presentation
                presentation.Close()
                powerpoint.Quit()
                
                # Read the PDF
                with open(tmp_pdf_path, 'rb') as pdf_file:
                    return pdf_file.read()
                    
            finally:
                # Clean up
                try:
                    os.unlink(tmp_pptx_path)
                    os.unlink(tmp_pdf_path)
                except:
                    pass
                    
        except Exception as e:
            logger.debug(f"Microsoft PowerPoint conversion failed: {e}")
            return None
    
    @staticmethod
    def _pptx_text_fallback(file_content: bytes) -> Optional[bytes]:
        """Fallback: Extract text from PPTX when LibreOffice is unavailable"""
        if not HAS_PPTX:
            return None
        
        try:
            prs = Presentation(io.BytesIO(file_content))
            pdf_buffer = io.BytesIO()
            doc = SimpleDocTemplate(pdf_buffer, pagesize=letter)
            styles = getSampleStyleSheet()
            story = []
            
            for slide_idx, slide in enumerate(prs.slides, 1):
                if slide_idx > 1:
                    story.append(PageBreak())
                story.append(Paragraph(f"<b>Slide {slide_idx}</b>", styles['Heading2']))
                story.append(Spacer(1, 0.1 * inch))
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        story.append(Paragraph(shape.text, styles['Normal']))
                        story.append(Spacer(1, 0.1 * inch))
            
            if not story:
                story.append(Paragraph("Empty presentation", styles['Normal']))
            
            doc.build(story)
            pdf_buffer.seek(0)
            return pdf_buffer.getvalue()
        except Exception as e:
            logger.error(f"Fallback PPTX conversion failed: {e}", exc_info=True)
            return None
    
    @staticmethod
    def _convert_with_libreoffice(file_content: bytes, file_ext: str) -> Optional[bytes]:
        """Convert files using LibreOffice in headless mode (works on Windows, Linux, Mac)"""
        try:
            # Check if LibreOffice is available
            soffice_cmd = None
            for cmd in ['soffice', 'libreoffice', r'C:\Program Files\LibreOffice\program\soffice.exe']:
                if shutil.which(cmd) or os.path.exists(cmd):
                    soffice_cmd = cmd
                    break
            
            if not soffice_cmd:
                logger.error("LibreOffice not found. Please install LibreOffice for document conversion.")
                return None
            
            # Create temp directory and files
            with tempfile.TemporaryDirectory() as temp_dir:
                input_file = os.path.join(temp_dir, f"input{file_ext}")
                
                # Write input file
                with open(input_file, 'wb') as f:
                    f.write(file_content)
                
                # Convert to PDF using LibreOffice
                cmd = [
                    soffice_cmd,
                    '--headless',
                    '--convert-to', 'pdf',
                    '--outdir', temp_dir,
                    input_file
                ]
                
                result = subprocess.run(cmd, capture_output=True, timeout=60, text=True)
                
                if result.returncode != 0:
                    logger.error(f"LibreOffice conversion failed: {result.stderr}")
                    return None
                
                # Read the output PDF
                output_file = os.path.join(temp_dir, "input.pdf")
                
                if not os.path.exists(output_file):
                    logger.error(f"PDF not created at {output_file}")
                    return None
                
                with open(output_file, 'rb') as f:
                    return f.read()
                    
        except subprocess.TimeoutExpired:
            logger.error("LibreOffice conversion timed out")
            return None
        except Exception as e:
            logger.error(f"Error in LibreOffice conversion: {e}", exc_info=True)
            return None
    
    @staticmethod
    def txt_to_pdf(file_content: bytes) -> Optional[bytes]:
        """Convert TXT to PDF"""
        try:
            text = file_content.decode('utf-8', errors='ignore')
            
            pdf_buffer = io.BytesIO()
            doc = SimpleDocTemplate(pdf_buffer, pagesize=letter)
            styles = getSampleStyleSheet()
            story = []
            
            for line in text.split('\n'):
                if line.strip():
                    story.append(Paragraph(line, styles['Normal']))
                else:
                    story.append(Spacer(1, 0.1 * inch))
            
            doc.build(story)
            pdf_buffer.seek(0)
            return pdf_buffer.getvalue()
            
        except Exception as e:
            logger.error(f"Error converting TXT: {e}", exc_info=True)
            return None
    
    @staticmethod
    def image_to_pdf(file_content: bytes) -> Optional[bytes]:
        """Convert Image to PDF"""
        try:
            img = PILImage.open(io.BytesIO(file_content))
            
            pdf_buffer = io.BytesIO()
            doc = SimpleDocTemplate(pdf_buffer, pagesize=letter)
            story = []
            
            # Resize if needed
            max_width = 6 * inch
            max_height = 9 * inch
            img_width, img_height = img.size
            ratio = min(max_width / img_width, max_height / img_height, 1.0)
            
            img_buffer = io.BytesIO()
            img.save(img_buffer, format='PNG')
            img_buffer.seek(0)
            
            from reportlab.platypus import Image as RLImage
            rl_img = RLImage(img_buffer, width=img_width * ratio, height=img_height * ratio)
            story.append(rl_img)
            
            doc.build(story)
            pdf_buffer.seek(0)
            return pdf_buffer.getvalue()
            
        except Exception as e:
            logger.error(f"Error converting image: {e}", exc_info=True)
            return None
    
    @staticmethod
    def convert(file_content: bytes, file_name: str) -> Tuple[Optional[bytes], str]:
        """
        Convert document to PDF
        Returns: (pdf_bytes, error_message)
        """
        ext = Path(file_name).suffix.lower().lstrip('.')
        
        try:
            if ext in ('docx', 'doc'):
                pdf = DocumentConverter.docx_to_pdf(file_content)
            elif ext in ('xlsx', 'xls'):
                pdf = DocumentConverter.xlsx_to_pdf(file_content)
            elif ext in ('pptx', 'ppt'):
                pdf = DocumentConverter.pptx_to_pdf(file_content)
            elif ext == 'txt':
                pdf = DocumentConverter.txt_to_pdf(file_content)
            elif ext in ('png', 'jpg', 'jpeg'):
                pdf = DocumentConverter.image_to_pdf(file_content)
            else:
                return None, f"Unsupported format: {ext}"
            
            if pdf is None:
                if ext == 'doc':
                    return None, "Old .doc format not supported. Please save as .docx"
                return None, f"Failed to convert {ext.upper()} file"
            
            return pdf, ""
            
        except Exception as e:
            logger.error(f"Conversion error: {e}", exc_info=True)
            return None, str(e)
